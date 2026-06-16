#!/usr/bin/env python3
"""
Generador de informe financiero mensual NLACE (.pptx)
Uso: python3 build_pptx.py <ruta_json> <ruta_salida.pptx>
"""
import json
import sys
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

# ── Identidad visual NLACE ──────────────────────────────────────────────────
BG          = RGBColor(0xEF, 0xEF, 0xEF)   # #EFEFEF canvas
TEXT        = RGBColor(0x0F, 0x10, 0x11)   # #0F1011 principal
BLUE        = RGBColor(0x58, 0x69, 0xF7)   # #5869F7 ingresos/ventas
ACCENT      = RGBColor(0xFC, 0x62, 0x4B)   # #FC624B costos/margen
GREEN       = RGBColor(0x22, 0xC5, 0x5E)   # #22C55E positivo
RED         = RGBColor(0xDC, 0x26, 0x26)   # #DC2626 negativo/gastos
GREY        = RGBColor(0x88, 0x88, 0x88)   # gris labels
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

FONT_TITLE  = "Space Grotesk"
FONT_BODY   = "Inter"

SLIDE_W     = Inches(13.333)
SLIDE_H     = Inches(7.5)

LOGO_PATH   = Path(__file__).parent / "assets" / "nlace-black.png"
LOGO_W      = Inches(1.4)
LOGO_H      = Inches(0.38)
LOGO_L      = Inches(0.5)
LOGO_T      = Inches(0.22)


# ── Helpers ──────────────────────────────────────────────────────────────────

def rgb_hex(color_str: str) -> RGBColor:
    h = color_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_bg(slide, color: RGBColor = BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def strip_all_shadows(slide):
    """Elimina sombras de todos los elementos del slide."""
    for shape in slide.shapes:
        try:
            sp_pr = shape.element.find(qn('p:spPr'))
            if sp_pr is None:
                sp_pr = getattr(shape.element, 'spPr', None)
            if sp_pr is None:
                continue
            for eff in list(sp_pr.findall(qn('a:effectLst'))):
                sp_pr.remove(eff)
            etree.SubElement(sp_pr, qn('a:effectLst'))
        except Exception:
            pass


def add_logo(slide, prs):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), LOGO_L, LOGO_T, LOGO_W, LOGO_H)


def add_footer(slide, page_num: int, total: int = 6):
    txBox = slide.shapes.add_textbox(
        Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.25)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page_num} / {total}"
    run.font.size = Pt(9)
    run.font.color.rgb = GREY
    run.font.name = FONT_BODY


def add_title(slide, title: str, subtitle: str = "", y: float = 0.55):
    left, top, width, height = Inches(0.5), Inches(y), Inches(12.3), Inches(0.7)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(28)
    run.font.name = FONT_TITLE
    run.font.color.rgb = TEXT

    if subtitle:
        left2, top2, width2, height2 = Inches(0.5), Inches(y + 0.65), Inches(12.3), Inches(0.35)
        txBox2 = slide.shapes.add_textbox(left2, top2, width2, height2)
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(13)
        run2.font.name = FONT_BODY
        run2.font.color.rgb = GREY


def add_divider(slide, y: float = 1.35):
    left, top, width, height = Inches(0.5), Inches(y), Inches(12.3), Inches(0.015)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = BLUE
    shape.line.fill.background()


def add_textbox(slide, text: str, left, top, width, height,
                font_name=FONT_BODY, font_size=11, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_kpi_block(slide, label: str, value: str, delta: str,
                  left, top, width=Inches(2.5), delta_color=None):
    if delta_color is None:
        delta_color = GREEN if "+" in delta else (RED if "-" in delta else GREY)

    # valor
    vbox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = vbox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = value
    run.font.name = FONT_TITLE
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = TEXT

    # delta
    dbox = slide.shapes.add_textbox(left, top + Inches(0.65), width, Inches(0.3))
    tf2 = dbox.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = delta
    run2.font.name = FONT_BODY
    run2.font.size = Pt(11)
    run2.font.color.rgb = delta_color

    # label
    lbox = slide.shapes.add_textbox(left, top + Inches(0.95), width, Inches(0.3))
    tf3 = lbox.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = label.upper()
    run3.font.name = FONT_BODY
    run3.font.size = Pt(9)
    run3.font.color.rgb = GREY


def format_clp(n, compact=False):
    if n is None:
        return "—"
    try:
        n = float(n)
    except (TypeError, ValueError):
        return "—"
    if compact:
        if abs(n) >= 1_000_000:
            return f"${n/1_000_000:.1f}M"
        if abs(n) >= 1_000:
            return f"${n/1_000:.0f}K"
    import locale
    try:
        locale.setlocale(locale.LC_ALL, '')
        return f"${locale.format_string('%,.0f', n)}"
    except Exception:
        return f"${n:,.0f}"


def format_pct(n):
    if n is None:
        return "—"
    try:
        return f"{float(n):.1f}%"
    except (TypeError, ValueError):
        return "—"


def delta_str(current, previous, pct=True):
    try:
        c, p = float(current or 0), float(previous or 0)
        if p == 0:
            return "+∞%" if c > 0 else "—"
        d = (c - p) / abs(p) * 100
        sign = "+" if d >= 0 else ""
        if pct:
            return f"{sign}{d:.1f}%"
        return f"{sign}{d:.0f}%"
    except (TypeError, ValueError):
        return "—"


def delta_color(current, previous, invert=False):
    try:
        c, p = float(current or 0), float(previous or 0)
        up = c >= p
        if invert:
            up = not up
        return GREEN if up else RED
    except (TypeError, ValueError):
        return GREY


# ── Slides ───────────────────────────────────────────────────────────────────

def slide_portada(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    add_logo(slide, prs)

    # Título grande centrado
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Informe Financiero Mensual"
    run.font.name = FONT_TITLE
    run.font.bold = True
    run.font.size = Pt(40)
    run.font.color.rgb = TEXT

    # Mes/año
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.85), Inches(10.3), Inches(0.7))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = data.get("periodo_label", "")
    run2.font.name = FONT_TITLE
    run2.font.size = Pt(24)
    run2.font.color.rgb = BLUE

    # Empresa
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(4.65), Inches(10.3), Inches(0.45))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "NLACE  ·  Preparado por CFO"
    run3.font.name = FONT_BODY
    run3.font.size = Pt(13)
    run3.font.color.rgb = GREY

    # Línea decorativa
    add_divider(slide, y=4.45)
    add_footer(slide, 1)
    strip_all_shadows(slide)


def slide_resumen_ejecutivo(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_logo(slide, prs)
    add_title(slide, "Resumen Ejecutivo", data.get("periodo_label", ""))
    add_divider(slide)
    add_footer(slide, 2)

    analisis = data.get("analisis", {})
    resumen = analisis.get("resumen", "")
    insights = analisis.get("insights", [])

    # Texto resumen (columna izquierda)
    add_textbox(slide, resumen,
                Inches(0.5), Inches(1.6), Inches(7.2), Inches(2.8),
                font_size=11, word_wrap=True)

    # Insights (columna derecha)
    if insights:
        header_box = slide.shapes.add_textbox(Inches(8.1), Inches(1.6), Inches(4.7), Inches(0.35))
        hf = header_box.text_frame
        hp = hf.paragraphs[0]
        hr = hp.add_run()
        hr.text = "PUNTOS CLAVE"
        hr.font.name = FONT_TITLE
        hr.font.bold = True
        hr.font.size = Pt(10)
        hr.font.color.rgb = BLUE

        for i, insight in enumerate(insights[:5]):
            y_pos = Inches(2.1 + i * 0.92)
            # bullet con color
            bul = slide.shapes.add_textbox(Inches(8.1), y_pos, Inches(0.25), Inches(0.3))
            bf = bul.text_frame
            bp = bf.paragraphs[0]
            br = bp.add_run()
            br.text = "●"
            br.font.color.rgb = BLUE
            br.font.size = Pt(9)

            txt = slide.shapes.add_textbox(Inches(8.45), y_pos, Inches(4.35), Inches(0.78))
            tf = txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = insight
            run.font.name = FONT_BODY
            run.font.size = Pt(10)
            run.font.color.rgb = TEXT

    strip_all_shadows(slide)


def slide_ingresos_margen(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_logo(slide, prs)
    add_title(slide, "Ingresos y Margen Operacional",
              f"Comparativo {data.get('mes_ant_label','')} vs {data.get('mes_label','')}")
    add_divider(slide)
    add_footer(slide, 3)

    m = data.get("metricas", {})
    ma = data.get("metricas_ant", {})

    # ── KPIs izquierda ──
    kpis = [
        ("Ingresos", m.get("ventas", 0), ma.get("ventas", 0), False),
        ("Costo de ventas", m.get("costos", 0), ma.get("costos", 0), True),
        ("Margen bruto", m.get("margen_bruto", 0), ma.get("margen_bruto", 0), False),
        ("Margen bruto %", None, None, False),
        ("Gastos operac.", m.get("gastos", 0), ma.get("gastos", 0), True),
        ("Resultado operac.", m.get("resultado_op", 0), ma.get("resultado_op", 0), False),
    ]
    for i, (label, val, val_ant, invert) in enumerate(kpis[:4]):
        row = i % 2
        col = i // 2
        lft = Inches(0.5 + col * 2.9)
        top = Inches(1.7 + row * 1.35)
        if label == "Margen bruto %":
            v_str = format_pct(m.get("margen_bruto_pct"))
            d_str = delta_str(m.get("margen_bruto_pct"), ma.get("margen_bruto_pct"))
            dc = delta_color(m.get("margen_bruto_pct"), ma.get("margen_bruto_pct"))
        else:
            v_str = format_clp(val, compact=True)
            d_str = delta_str(val, val_ant)
            dc = delta_color(val, val_ant, invert=invert)
        add_kpi_block(slide, label, v_str, d_str, lft, top, delta_color=dc)

    # ── Gráfico barras derecha ──
    chart_data = CategoryChartData()
    chart_data.categories = [data.get("mes_ant_label", "M-1"), data.get("mes_label", "M")]
    ventas_vals = [ma.get("ventas", 0) or 0, m.get("ventas", 0) or 0]
    margen_vals = [ma.get("margen_bruto", 0) or 0, m.get("margen_bruto", 0) or 0]
    chart_data.add_series("Ingresos", ventas_vals)
    chart_data.add_series("Margen bruto", margen_vals)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7.0), Inches(1.55), Inches(5.9), Inches(4.6),
        chart_data,
    ).chart

    # Colores planos sin gradiente
    series_colors = [BLUE, ACCENT]
    for idx, series in enumerate(chart.series):
        sp_pr = series.format.fill
        sp_pr.solid()
        sp_pr.fore_color.rgb = series_colors[idx]

    chart.has_legend = True
    chart.legend.position = 2  # bottom
    chart.plots[0].gap_width = 80

    # Fondo del chart transparente (API python-pptx 1.x)
    try:
        chart.plot_area.format.fill.background()
    except Exception:
        pass

    strip_all_shadows(slide)


def slide_gastos(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_logo(slide, prs)
    add_title(slide, "Gastos Operacionales por Clasificación",
              f"Comparativo {data.get('mes_ant_label','')} vs {data.get('mes_label','')}")
    add_divider(slide)
    add_footer(slide, 4)

    gastos = data.get("gastos_clasif", [])  # [{clasificacion, mes_ant, mes_actual}]
    lectura = data.get("analisis", {}).get("lectura_gastos", "")

    # Texto lectura (izquierda)
    add_textbox(slide, lectura,
                Inches(0.5), Inches(1.7), Inches(5.5), Inches(4.0),
                font_size=11, word_wrap=True)

    # Tabla (derecha)
    if gastos:
        headers = ["Clasificación", data.get("mes_ant_label", "M-1"), data.get("mes_label", "M"), "Δ%"]
        col_widths = [Inches(2.4), Inches(1.3), Inches(1.3), Inches(0.9)]
        start_x = Inches(6.4)
        start_y = Inches(1.6)
        row_h = Inches(0.4)

        # Header row
        for j, (hdr, w) in enumerate(zip(headers, col_widths)):
            x = start_x + sum(col_widths[:j])
            hbox = slide.shapes.add_textbox(x, start_y, w, row_h)
            tf = hbox.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = hdr
            run.font.name = FONT_BODY
            run.font.bold = True
            run.font.size = Pt(9)
            run.font.color.rgb = BLUE

        # Divider bajo header
        div = slide.shapes.add_shape(1, start_x, start_y + row_h,
                                      sum(col_widths), Inches(0.012))
        div.fill.solid()
        div.fill.fore_color.rgb = GREY
        div.line.fill.background()

        for i, g in enumerate(gastos[:12]):
            y = start_y + row_h + Inches(0.015) + i * Inches(0.38)
            row_vals = [
                g.get("clasificacion", ""),
                format_clp(g.get("mes_ant", 0), compact=True),
                format_clp(g.get("mes_actual", 0), compact=True),
                delta_str(g.get("mes_actual", 0), g.get("mes_ant", 0)),
            ]
            for j, (val, w) in enumerate(zip(row_vals, col_widths)):
                x = start_x + sum(col_widths[:j])
                vbox = slide.shapes.add_textbox(x, y, w, Inches(0.35))
                tf = vbox.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
                run = p.add_run()
                run.text = val
                run.font.name = FONT_BODY
                run.font.size = Pt(9)
                if j == 3:
                    try:
                        d = float(g.get("mes_actual", 0) or 0)
                        a = float(g.get("mes_ant", 0) or 0)
                        run.font.color.rgb = RED if d > a else GREEN
                    except Exception:
                        run.font.color.rgb = GREY
                else:
                    run.font.color.rgb = TEXT

    strip_all_shadows(slide)


def slide_cobranza(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_logo(slide, prs)
    add_title(slide, "Estado de Cobranza (CxC)", "Facturas emitidas pendientes de pago")
    add_divider(slide)
    add_footer(slide, 5)

    cxc = data.get("cobranza", {})
    monto_pend = cxc.get("monto_pendiente", 0)
    num_facturas = cxc.get("num_facturas", 0)
    tasa_cobro = cxc.get("tasa_cobro_pct", None)
    dso = cxc.get("dso_dias", None)
    top_deudores = cxc.get("top_deudores", [])
    lectura = data.get("analisis", {}).get("lectura_cobranza", "")

    # KPIs izquierda
    add_kpi_block(slide, "Monto pendiente",
                  format_clp(monto_pend, compact=True), "", Inches(0.5), Inches(1.7),
                  delta_color=RED if (monto_pend or 0) > 0 else GREEN)

    add_kpi_block(slide, "Facturas sin cobrar",
                  str(num_facturas), "", Inches(0.5), Inches(3.0), delta_color=GREY)

    if tasa_cobro is not None:
        tc_color = GREEN if float(tasa_cobro) >= 80 else (ACCENT if float(tasa_cobro) >= 60 else RED)
        add_kpi_block(slide, "Tasa de cobro",
                      format_pct(tasa_cobro), "", Inches(3.1), Inches(1.7), delta_color=tc_color)

    if dso is not None:
        dso_color = GREEN if float(dso) <= 30 else (ACCENT if float(dso) <= 45 else RED)
        add_kpi_block(slide, "DSO promedio",
                      f"{float(dso):.0f} días", "", Inches(3.1), Inches(3.0), delta_color=dso_color)

    # Texto lectura
    add_textbox(slide, lectura,
                Inches(0.5), Inches(4.3), Inches(5.5), Inches(1.8),
                font_size=10, word_wrap=True)

    # Top deudores (derecha)
    if top_deudores:
        hdr = slide.shapes.add_textbox(Inches(6.5), Inches(1.7), Inches(6.4), Inches(0.35))
        ht = hdr.text_frame
        hp = ht.paragraphs[0]
        hr = hp.add_run()
        hr.text = "TOP DEUDORES"
        hr.font.name = FONT_TITLE
        hr.font.bold = True
        hr.font.size = Pt(10)
        hr.font.color.rgb = BLUE

        for i, d in enumerate(top_deudores[:6]):
            y = Inches(2.15 + i * 0.75)
            # nombre
            nb = slide.shapes.add_textbox(Inches(6.5), y, Inches(4.0), Inches(0.3))
            ntf = nb.text_frame
            np_ = ntf.paragraphs[0]
            nr = np_.add_run()
            nr.text = d.get("cliente", "")[:40]
            nr.font.name = FONT_BODY
            nr.font.bold = True
            nr.font.size = Pt(10)
            nr.font.color.rgb = TEXT
            # monto
            mb = slide.shapes.add_textbox(Inches(10.7), y, Inches(2.1), Inches(0.3))
            mtf = mb.text_frame
            mp = mtf.paragraphs[0]
            mp.alignment = PP_ALIGN.RIGHT
            mr = mp.add_run()
            mr.text = format_clp(d.get("monto", 0), compact=True)
            mr.font.name = FONT_BODY
            mr.font.size = Pt(10)
            mr.font.color.rgb = ACCENT
            # días mora
            if d.get("dias_mora") is not None:
                db = slide.shapes.add_textbox(Inches(6.5), y + Inches(0.3), Inches(4.0), Inches(0.25))
                dtf = db.text_frame
                dpp = dtf.paragraphs[0]
                drn = dpp.add_run()
                dias = int(d.get("dias_mora", 0))
                drn.text = f"{dias} días de mora"
                drn.font.name = FONT_BODY
                drn.font.size = Pt(8.5)
                drn.font.color.rgb = RED if dias > 30 else GREY

    strip_all_shadows(slide)


def slide_resultado_neto(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_logo(slide, prs)
    add_title(slide, "Resultado Neto",
              f"Estado de resultado — {data.get('mes_label','')}")
    add_divider(slide)
    add_footer(slide, 6)

    m = data.get("metricas", {})
    ma = data.get("metricas_ant", {})
    semaforo = data.get("semaforo", "verde")

    # Mapa del Estado de Resultado (izquierda)
    lineas = [
        ("Ingresos",            m.get("ventas", 0),          ma.get("ventas", 0),          False),
        ("– Costo de ventas",   m.get("costos", 0),          ma.get("costos", 0),          True),
        ("= Margen bruto",      m.get("margen_bruto", 0),    ma.get("margen_bruto", 0),    False),
        ("– Gastos operac.",    m.get("gastos", 0),          ma.get("gastos", 0),          True),
        ("= Resultado operac.", m.get("resultado_op", 0),    ma.get("resultado_op", 0),    False),
        ("– Rem. directores",   m.get("rem_directores", 0),  ma.get("rem_directores", 0),  True),
        ("= EBITDA",            m.get("ebitda", 0),          ma.get("ebitda", 0),          False),
    ]
    bold_rows = {2, 4, 6}

    y0 = Inches(1.65)
    for i, (label, val, val_ant, invert) in enumerate(lineas):
        y = y0 + Inches(i * 0.68)
        is_bold = i in bold_rows

        # separador antes de líneas importantes
        if i in {2, 4, 6}:
            sep = slide.shapes.add_shape(1, Inches(0.5), y - Inches(0.08),
                                          Inches(6.8), Inches(0.008))
            sep.fill.solid()
            sep.fill.fore_color.rgb = GREY
            sep.line.fill.background()

        lbox = slide.shapes.add_textbox(Inches(0.5), y, Inches(3.5), Inches(0.55))
        ltf = lbox.text_frame
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = label
        lr.font.name = FONT_BODY
        lr.font.bold = is_bold
        lr.font.size = Pt(11 if is_bold else 10)
        lr.font.color.rgb = TEXT

        vbox = slide.shapes.add_textbox(Inches(4.2), y, Inches(1.6), Inches(0.55))
        vtf = vbox.text_frame
        vp = vtf.paragraphs[0]
        vp.alignment = PP_ALIGN.RIGHT
        vr = vp.add_run()
        vr.text = format_clp(val, compact=True)
        vr.font.name = FONT_BODY
        vr.font.bold = is_bold
        vr.font.size = Pt(11 if is_bold else 10)
        vr.font.color.rgb = TEXT

        # delta %
        d_str = delta_str(val, val_ant)
        dc = delta_color(val, val_ant, invert=invert)
        dbox = slide.shapes.add_textbox(Inches(5.9), y, Inches(0.9), Inches(0.55))
        dtf = dbox.text_frame
        dp = dtf.paragraphs[0]
        dp.alignment = PP_ALIGN.RIGHT
        dr = dp.add_run()
        dr.text = d_str
        dr.font.name = FONT_BODY
        dr.font.size = Pt(9)
        dr.font.color.rgb = dc

    # ── Semáforo + KPIs principales (derecha) ──
    # Círculo semáforo
    sem_colors = {"verde": GREEN, "amarillo": ACCENT, "rojo": RED}
    sem_col = sem_colors.get(semaforo, GREY)
    circ = slide.shapes.add_shape(
        9,  # oval
        Inches(8.2), Inches(2.0), Inches(1.5), Inches(1.5)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = sem_col
    circ.line.fill.background()

    # Etiqueta semáforo
    sl = slide.shapes.add_textbox(Inches(8.0), Inches(3.65), Inches(2.0), Inches(0.35))
    stf = sl.text_frame
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = semaforo.upper()
    sr.font.name = FONT_TITLE
    sr.font.bold = True
    sr.font.size = Pt(13)
    sr.font.color.rgb = sem_col

    # KPIs resultado operac. + EBITDA
    res_op = m.get("resultado_op", 0)
    ebitda = m.get("ebitda", 0)
    for label, val, val_ant, lx in [
        ("Resultado operac.", res_op, ma.get("resultado_op", 0), Inches(10.0)),
        ("EBITDA", ebitda, ma.get("ebitda", 0), Inches(11.7)),
    ]:
        col_v = GREEN if (val or 0) >= 0 else RED
        add_kpi_block(slide, label, format_clp(val, compact=True),
                      delta_str(val, val_ant),
                      lx, Inches(2.0),
                      width=Inches(1.5),
                      delta_color=delta_color(val, val_ant))

    strip_all_shadows(slide)


# ── Main ─────────────────────────────────────────────────────────────────────

def build(json_path: str, out_path: str):
    with open(json_path, encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_portada(prs, data)
    slide_resumen_ejecutivo(prs, data)
    slide_ingresos_margen(prs, data)
    slide_gastos(prs, data)
    slide_cobranza(prs, data)
    slide_resultado_neto(prs, data)

    prs.save(out_path)
    print(f"✓ Informe guardado en: {out_path}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Uso: python3 build_pptx.py <datos.json> <salida.pptx>")
        sys.exit(1)
    build(sys.argv[1], sys.argv[2])
